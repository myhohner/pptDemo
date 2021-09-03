from pptx import Presentation
from pptx.util import Pt  #pt 像素单位
import pandas as pd
import os
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
path=os.path.dirname(__file__)

class WritePowerPoint:
    def __init__(self, ppt_name, input_excel, title_cover):
        self.ppt_name = ppt_name
        self.input_excel = input_excel
        self.title_cover = self.title_per_page = title_cover
        # 创建空白演示文稿
        self.prs = Presentation()
          
    def run(self):
        slide_layout=self.prs.slide_layouts[5]
        slide=self.prs.slides.add_slide(slide_layout)
        self.shapes=slide.shapes

        title_shape=self.shapes.title   #主标题
        title_shape.text=self.title_cover #主标题内容
        title_shape.left, title_shape.top = Pt(28*2.5), Pt(28*0.37)
        title_shape.width, title_shape.height = Pt(28*20), Pt(28*1.11)
        p=title_shape.text_frame.paragraphs[0]
        p.font.size=Pt(20) #标题大小
        p.font.name = '微软雅黑'
        p.alignment =PP_ALIGN.LEFT #左对齐
        
        '''
        sub_title_shape=slide.placeholders[1]  #副标题
        sub_title_shape.text=self.subtitle_cover
        sub_title_shape.left,sub_title_shape.top=Pt(28*2.5), Pt(28*8)
        sub_title_shape.width, sub_title_shape.height = Pt(28*20), Pt(28*1.11)
        p=sub_title_shape.text_frame.paragraphs[0]
        p.font.size=Pt(20) #标题大小
        p.font.name = '微软雅黑'
        p.alignment =PP_ALIGN.LEFT #左对齐
        '''
        
        tf_heat=self.subtitle_font(2,'采暖行业')
        tf_house=self.subtitle_font(7.5,'房地产行业')
       
        #读取excel
        df=pd.read_excel(self.input_excel,sheet_name = '行业')
        df1=df[df['摘要'].notnull()]
        df_heat=df1[df1['内容分类'].str.contains('采暖')].astype(str)
        df_house=df1[df1['内容分类'].str.contains('房地产')].astype(str)
        self.add_texts(df_heat,tf_heat)
        self.add_texts(df_house,tf_house)
        
        self.prs.save(self.ppt_name)
   
    def subtitle_font(self,top_Pt,text):
        left, top, width, height = Pt(28*1), Pt(top_Pt*28), Pt(28*20), Pt(28*6)
        text_box = self.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        
        p=tf.paragraphs[0]
        #p.text='房地产行业'
        p.text=text
        p.font.size=Pt(12) 
        p.font.name = '微软雅黑'
        p.font.bold= True
        tf.add_paragraph()
        
        return tf
   
    def add_texts(self,df,tf):
        num=8
        n=num if num<=len(df) else len(df)
        for i in range(n):
            text_list=df.iloc[i,:].tolist()
            self.add_paragraph_texts(text_list,tf,i)
            
    def add_paragraph_texts(self,text_list,tf,i):
        title=text_list[1]
        #summary=text_list[3]     
        link=text_list[13]
        media=text_list[6]
        p = tf.paragraphs[i+1]
        p.text=title+'-'
        p.font.size=Pt(11) 
        p.font.name = '微软雅黑'
        p.line_spacing = 1.5
        run = p.add_run() #增加超链接
        run.text =media
        run.hyperlink.address =link
        tf.add_paragraph()
        
if __name__ == '__main__':
    wpt = WritePowerPoint(path+'/news.pptx', path+r'/数据.xls', '行业新闻分析-主要新闻展示')
    wpt.run()